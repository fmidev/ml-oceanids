#!/usr/bin/env python3
"""
Script to generate PowerPoint slides from SHAP analysis results.
Creates one slide per harbor per predictor with SHAP visualizations.
"""

import os
import sys
import json
import argparse
from datetime import datetime
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# List of predictors to process
PREDICTORS = [
    'TA_PT24H_MAX',
    'TA_PT24H_MIN',
    'WS_PT24H_AVG',
    'RH_PT24H_AVG',
    'TP_PT24H_ACC',
    'WG_PT24H_MAX'
]

# Mapping of predictor codes to human-readable descriptions
PREDICTOR_DESCRIPTIONS = {
    'TA_PT24H_MAX': 'Maximum Temperature',
    'TA_PT24H_MIN': 'Minimum Temperature',
    'WS_PT24H_AVG': 'Average Wind Speed',
    'RH_PT24H_AVG': 'Average Relative Humidity',
    'TP_PT24H_ACC': 'Total Precipitation',
    'WG_PT24H_MAX': 'Maximum Wind Gust'
}

def get_harbors_from_config(config_path):
    """Get list of harbors from the configuration file."""
    with open(config_path, 'r') as f:
        harbors_config = json.load(f)
    return list(harbors_config.keys())

def get_metrics_for_harbor_predictor(harbor, predictor, metrics_dir):
    """Load metrics data for a specific harbor and predictor."""
    metrics_file = f"{metrics_dir}/{harbor}-metrics.csv"
    
    if not os.path.exists(metrics_file):
        print(f"Metrics file not found: {metrics_file}")
        return None
    
    try:
        metrics_df = pd.read_csv(metrics_file)
        # Filter for the specific predictor
        predictor_metrics = metrics_df[metrics_df['Prediction'] == predictor]
        
        if predictor_metrics.empty:
            print(f"No metrics found for {harbor} - {predictor}")
            return None
            
        # Return the first (and should be only) matching row
        return predictor_metrics.iloc[0]
    except Exception as e:
        print(f"Error loading metrics for {harbor} - {predictor}: {e}")
        return None

def get_nan_percentage(harbor_name, predictor=None):
    """
    Calculate the percentage of NaN values in the training data.
    If predictor is specified, calculate NaN percentage for that specific predictor.
    """
    training_file = f'/home/ubuntu/data/ML/training-data/OCEANIDS/{harbor_name}/training_data_oceanids_{harbor_name}-sf-addpreds.csv.gz'
    
    try:
        # Check if file exists
        if not os.path.exists(training_file):
            print(f"Training data for {harbor_name} not found.")
            return None
            
        print(f"Found training data for {harbor_name}.")
        # Load the training data
        df = pd.read_csv(training_file, compression='gzip')
        
        # First drop columns that are all NaN
        df = df.dropna(axis=1, how='all')
        
        # If specific predictor requested, calculate NaN percentage for just that column
        if predictor:
            # Look for the exact predictor column
            if predictor in df.columns:
                pred_col = predictor
            else:
                # Try alternative column naming patterns if exact match not found
                possible_cols = [col for col in df.columns if predictor in col]
                if possible_cols:
                    pred_col = possible_cols[0]
                else:
                    return None
            
            # Calculate percentage of NaNs for this specific predictor
            total_rows = df.shape[0]
            nan_count = df[pred_col].isna().sum()
            
            if total_rows > 0:
                nan_percentage = round((nan_count / total_rows * 100), 1)
                print(f"{harbor_name} - {predictor}: Missing data = {nan_percentage}%")
                return nan_percentage
        else:
            # Calculate overall NaN percentage
            total_rows = df.shape[0]
            df_no_nans = df.dropna(axis=0, how='any')
            remaining_rows = df_no_nans.shape[0]
            
            if total_rows > 0:
                nan_percentage = round(100 - (remaining_rows / total_rows * 100), 1)
                print(f"{harbor_name} - overall: Missing data = {nan_percentage}%")
                return nan_percentage
        
        return 0
    except Exception as e:
        print(f"Error calculating NaN percentage for {harbor_name}: {e}")
        return None

def add_section_slide(prs, title):
    """Add a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # Section header layout
    title_shape = slide.shapes.title
    title_shape.text = title
    title_shape.text_frame.paragraphs[0].font.size = Pt(40)
    title_shape.text_frame.paragraphs[0].font.bold = True
    
    return slide

def add_harbor_predictor_slide(prs, harbor_name, pred, res_dir, plots_dir, metrics_dir, harbor_config):
    """Add slide for specific harbor and predictor with 16:9 aspect ratio."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # Title and content layout
    
    # Add title with specified styling
    title_shape = slide.shapes.title
    title_shape.text = f"{harbor_name.capitalize()}: {PREDICTOR_DESCRIPTIONS.get(pred, pred)}"
    
    # Style the title - font, color, size
    title_frame = title_shape.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.font.name = "Univa Nova"
    title_para.font.size = Pt(24)
    title_para.font.color.rgb = RGBColor.from_string('2E7E7C')
    title_frame.word_wrap = True
    title_para.alignment = PP_ALIGN.CENTER
    
    # Define image paths for variant SHAP images
    image_prefix = f"{harbor_name}_{pred}_xgb_era5_oceanids-QE"
    grid_comparison_path = f"{res_dir}grid_point_beeswarm_{image_prefix}.png"
    variable_beeswarm_path = f"{res_dir}variable_beeswarm_{image_prefix}.png"
    
    # Define path for harbor plot - use the correct filename pattern
    harbor_plot_path = f"{res_dir}{harbor_name}_training-locs.png"
    
    # Add variable beeswarm plot (perfectly aligned with top right)
    if os.path.exists(variable_beeswarm_path):
        width = Inches(5.2)  # Slightly smaller
        left = Inches(13.333) - width  # Align with right edge of slide
        top = Inches(0.75)  # Align perfectly with top
        var_beeswarm = slide.shapes.add_picture(variable_beeswarm_path, left, top, width=width)
    else:
        # Add placeholder text if image not found
        txBox = slide.shapes.add_textbox(Inches(8.0), Inches(3.0), Inches(5.0), Inches(1.0))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Variable beeswarm plot not available"
        p.font.italic = True
        p.font.size = Pt(14)
    
    # Add central text area with training information
    txBox = slide.shapes.add_textbox(Inches(4.5), Inches(1.8), Inches(3.5), Inches(3.0))
    tf = txBox.text_frame
    
    # Add header for data information
    p = tf.add_paragraph()
    p.text = "Data Information"
    p.font.bold = True
    p.font.size = Pt(18)
    
    # Add training timeframe from harbors config
    if harbor_config and 'start' in harbor_config and 'end' in harbor_config:
        # Format dates in dd.mm.yyyy format
        start_date_iso = harbor_config['start'].split('T')[0]  # Extract date part
        end_date_iso = harbor_config['end'].split('T')[0]
        
        # Convert to dd.mm.yyyy format
        try:
            # Try to parse with format '%Y%m%d' (20050101)
            start_dt = datetime.strptime(start_date_iso, '%Y%m%d')
            end_dt = datetime.strptime(end_date_iso, '%Y%m%d')
            start_date = start_dt.strftime('%d.%m.%Y')
            end_date = end_dt.strftime('%d.%m.%Y')
        except ValueError:
            try:
                # Fallback to ISO format '%Y-%m-%d' (2005-01-01)
                start_dt = datetime.strptime(start_date_iso, '%Y-%m-%d')
                end_dt = datetime.strptime(end_date_iso, '%Y-%m-%d')
                start_date = start_dt.strftime('%d.%m.%Y')
                end_date = end_dt.strftime('%d.%m.%Y')
            except ValueError:
                # Use original format if both parsing attempts fail
                start_date = start_date_iso
                end_date = end_date_iso
        
        p = tf.add_paragraph()
        p.text = f"Training period: {start_date} to {end_date}"
        p.font.size = Pt(14)
        p.level = 1
        
    # Calculate and add NaN percentage for this specific predictor
    nan_percentage = get_nan_percentage(harbor_name, pred)
    p = tf.add_paragraph()
    if nan_percentage is not None:
        p.text = f"Missing data: {nan_percentage}%"
    else:
        p.text = "Missing data: N/A"
    p.font.size = Pt(14)
    p.level = 1
    
    # Load and add metrics data if available
    metrics = get_metrics_for_harbor_predictor(harbor_name, pred, metrics_dir)
    if metrics is not None:
        # Add a metrics section
        p = tf.add_paragraph()
        p.text = "Performance Metrics:"
        p.font.bold = True
        p.font.size = Pt(16)
        
        # Add each metric as a bullet point
        metric_points = [
            f"Mean Error: {metrics['Mean_Error']:.3f}",
            f"MAE: {metrics['MAE']:.3f}",
            f"R²: {metrics['R2']:.3f}",
            f"Explained Variance: {metrics['Explained_Variance']:.3f}"
        ]
        
        for metric in metric_points:
            p = tf.add_paragraph()
            p.text = metric
            p.font.size = Pt(14)
            p.level = 1
    
    # Add grid point beeswarm plot (positioned under text)
    if os.path.exists(grid_comparison_path):
        left = Inches(3.8)  # Slightly more to the left
        top = Inches(4.8)   # Under the text
        width = Inches(4.5)  # Same width
        slide.shapes.add_picture(grid_comparison_path, left, top, width=width)
    else:
        # Add placeholder text if image not found
        txBox = slide.shapes.add_textbox(Inches(4.5), Inches(4.8), Inches(4.5), Inches(1.0))
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Grid point comparison not available"
        p.font.italic = True
        p.font.size = Pt(14)
    
    # Add harbor plot (better aligned with bottom of slide)
    if os.path.exists(harbor_plot_path):
        width = Inches(3.0)  # Keep same size
        left = Inches(0.8)  # Same horizontal position
        
        # Calculate better position to align with bottom
        # Using 7.5 (slide height) - 0.9 (bottom margin) - estimated height
        estimated_height = width * 0.75  # Estimate height based on width
        top = Inches(7.5 - 0.9 - estimated_height)
        
        slide.shapes.add_picture(harbor_plot_path, left, top, width=width)
    else:
        # Try alternative filenames if the standard one isn't found
        alternative_paths = [
            f"{res_dir}{harbor_name}_plot.png",
            f"{res_dir}{harbor_name}.png",
            f"{res_dir}map_{harbor_name}.png",
            f"{plots_dir}/{harbor_name}_map.png"
        ]
        
        found = False
        for alt_path in alternative_paths:
            if os.path.exists(alt_path):
                width = Inches(3.0)
                left = Inches(0.8)
                
                # Calculate top position using the same calculation as above
                estimated_height = width * 0.75
                top = Inches(7.5 - 0.9 - estimated_height)
                
                slide.shapes.add_picture(alt_path, left, top, width=width)
                found = True
                break
        
        if not found:
            # Add placeholder text if image not found
            txBox = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(3.0), Inches(1.0))
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = "Harbor plot not available"
            p.font.italic = True
            p.font.size = Pt(14)
    
    return slide

def parse_args():
    """Parse command line arguments."""
    parser = argparse.ArgumentParser(description='Create PowerPoint slides from SHAP analysis results')
    parser.add_argument('--output', '-o', help='Output file path', default=None)
    return parser.parse_args()

def create_presentation(args):
    """Create PowerPoint presentation from SHAP analysis results for all harbors and predictors."""
    # Get list of harbors from config
    harbors_config_path = '/home/ubuntu/ml-oceanids/bin/harbors_config.json'
    with open(harbors_config_path, 'r') as f:
        harbors_config = json.load(f)
    harbors = list(harbors_config.keys())
    
    # Define paths
    base_results_dir = '/home/ubuntu/data/ML/results/OCEANIDS/'
    plots_dir = f'{base_results_dir}'
    metrics_dir = f'{base_results_dir}metrics'
    
    # Define output file path
    if args.output:
        output_path = args.output
    else:
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = f"{base_results_dir}shap_presentation_all_harbors.pptx"
    
    # Create presentation
    prs = Presentation()
    
    # Set slide dimensions for 16:9 aspect ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Print summarized NaN percentages for all harbors and predictors
    print("\nCalculating NaN percentages for all harbors and predictors...")
    nan_percentages = {}
    for harbor in harbors:
        nan_percentages[harbor] = {}
        for pred in PREDICTORS:
            nan_percentages[harbor][pred] = get_nan_percentage(harbor, pred)
    
    # Process each harbor
    slides_created = 0
    for harbor in harbors:
        harbor_slides = 0
        res_dir = f'{base_results_dir}{harbor}/'
        
        # Process each predictor for this harbor
        for pred in PREDICTORS:
            # Check if variant SHAP analysis exists for this harbor-predictor combination
            image_prefix = f"{harbor}_{pred}_xgb_era5_oceanids-QE"
            grid_comparison_path = f"{res_dir}grid_point_comparison_{image_prefix}.png"
            
            if os.path.exists(grid_comparison_path):
                # Add a section slide for this harbor if this is the first predictor we found
                if harbor_slides == 0:
                    add_section_slide(prs, f"Harbor: {harbor.capitalize()}")
                
                print(f"Processing {harbor} - {pred}")
                add_harbor_predictor_slide(prs, harbor, pred, res_dir, plots_dir, metrics_dir, harbors_config.get(harbor))
                harbor_slides += 1
                slides_created += 1
            else:
                print(f"Skipping {harbor} - {pred}: Variant SHAP images not found")
    
    # Save presentation
    prs.save(output_path)
    print(f"Presentation created: {output_path}")
    print(f"Contains {slides_created} slides for {len(harbors)} harbors")
    return output_path


if __name__ == "__main__":
    args = parse_args()
    output_path = create_presentation(args)
