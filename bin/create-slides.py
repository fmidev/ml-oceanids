from pptx import Presentation
from pptx.util import Inches
import os
import glob

def create_slide_deck():
    # Initialize presentation
    prs = Presentation()
    
    # Define directories containing images
    image_dirs = [
        '/home/ubuntu/data/ML/results/OCEANIDS/',
        '/home/ubuntu/ml-harvesterseasons/',
        '/home/ubuntu/ml-oceanids/'
    ]
    
    # Supported image extensions
    image_extensions = ['*.jpg', '*.jpeg', '*.png']
    
    # Find all relevant images
    images = []
    for directory in image_dirs:
        for ext in image_extensions:
            images.extend(glob.glob(os.path.join(directory, ext)))
    
    # Create title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = title_slide.shapes.title
    subtitle = title_slide.placeholders[1]
    title.text = "OCEANIDS Analysis Results"
    subtitle.text = "Generated from ML analysis scripts"
    
    # Add images to slides
    for image_path in sorted(images):
        # Create a new slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Layout with title and content
        
        # Add title (using filename without extension)
        title = slide.shapes.title
        title.text = os.path.splitext(os.path.basename(image_path))[0]
        
        # Add image
        # Calculate image size while maintaining aspect ratio
        img_width = Inches(8)  # Max width
        left = Inches(1)
        top = Inches(2)
        
        slide.shapes.add_picture(image_path, left, top, width=img_width)
    
    # Save the presentation
    output_path = '/home/ubuntu/ml-oceanids/analysis_results.pptx'
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")

if __name__ == "__main__":
    create_slide_deck()
